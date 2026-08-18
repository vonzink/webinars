#!/usr/bin/env python3
"""
Generate The Homebuyer's Playbook as a PowerPoint (.pptx).

Ridgeline design system, 16:9. 16 main slides + 30 popout slides.
Cards on main slides jump to their popout; each popout has a "< Back" that
returns to its parent slide. Content mirrors deck/content/*.js.

Run:  python3 build_pptx.py   ->  ../Homebuyers-Playbook.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re, os

# ---- palette ---------------------------------------------------------------
FOREST=RGBColor(0x0C,0x33,0x35); FOREST_MID=RGBColor(0x14,0x49,0x4B)
SAGE=RGBColor(0x4b,0x7b,0x4d); GREEN=RGBColor(0x8c,0xc6,0x3E)
CHARCOAL=RGBColor(0x40,0x40,0x41); MIST=RGBColor(0xF5,0xF7,0xF4); WHITE=RGBColor(0xFF,0xFF,0xFF)
RED=RGBColor(0xC0,0x39,0x2B); DO_GREEN=RGBColor(0x3E,0x8E,0x41)
BODY_DARK=RGBColor(0xD6,0xDE,0xDA); META_DARK=RGBColor(0xA9,0xB6,0xB1); META_LIGHT=RGBColor(0x6A,0x73,0x70)
BANDS=[GREEN,RGBColor(0x6F,0xA8,0x47),SAGE,FOREST_MID,FOREST]
PAY=[FOREST,FOREST_MID,SAGE,RGBColor(0x8F,0xA6,0x9B),GREEN,CHARCOAL]

DISP="Montserrat"; BODY="Open Sans"
def IN(px): return Inches(px*13.333/1920.0)
def PT(px): return Pt(px*0.5)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
PADX=IN(96); PADY=IN(90); CW=Inches(13.333); CH=Inches(7.5)

# ---- compliance / presenter ----
COMPANY_L1="Mountain State Financial Group, LLC · NMLS# 1314257 · msfg.us"
COMPANY_L2="Licensed in CO, IN, MI, MN, ND, SD, TX · Equal Housing Lender"
SETH=dict(name="Seth Angell",title="Executive VP · NMLS# 912881",
          phone="303-883-8519",email="Seth.angell@msfg.us")
APPLY_URL="https://www.blink.mortgage/app/signup/p/mountainstatefinancialgroupllc/sethangell"
LOGO=os.path.join(os.path.dirname(__file__),"assets/brand/logo-horizontal.png")
QR=os.path.join(os.path.dirname(__file__),"assets/brand/qr-seth.png")
PORTRAIT=os.path.join(os.path.dirname(__file__),"assets/portraits/seth-angell.png")

# ---------------------------------------------------------------------------
def bg(slide,color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=color

def rect(slide,l,t,w,h,fill=None,line=None,shape=MSO_SHAPE.RECTANGLE):
    sp=slide.shapes.add_shape(shape,l,t,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    return sp

RICH=re.compile(r'(<strong>.*?</strong>|<em>.*?</em>)')
def add_rich(p,html,size,color,font=BODY,bold=False):
    for part in RICH.split(html):
        if not part: continue
        b=bold; txt=part
        if part.startswith('<strong>'): b=True; txt=part[8:-9]
        elif part.startswith('<em>'): b=True; txt=part[4:-5]
        txt=re.sub('<[^>]+>','',txt)
        r=p.add_run(); r.text=txt
        r.font.name=font; r.font.size=size; r.font.bold=b; r.font.color.rgb=color

def textbox(slide,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def line0(tf):  # first paragraph helper
    return tf.paragraphs[0]

def para(tf,space_before=0,space_after=6,align=PP_ALIGN.LEFT):
    p=tf.add_paragraph(); p.space_before=Pt(space_before); p.space_after=Pt(space_after); p.alignment=align
    return p

def set_text(p,txt,size,color,font=DISP,bold=True,align=PP_ALIGN.LEFT):
    p.alignment=align; r=p.add_run(); r.text=txt
    r.font.name=font; r.font.size=size; r.font.bold=bold; r.font.color.rgb=color

# ---- shared furniture ------------------------------------------------------
def eyebrow(slide,txt,dark):
    tf=textbox(slide,PADX,PADY,IN(1400),IN(40))
    set_text(line0(tf),txt.upper(),PT(23),GREEN if dark else SAGE,DISP,True)
    line0(tf).runs[0].font._rPr.set('spc','400')

def headline(slide,txt,dark,y=None,size=72,color=None):
    y=y if y is not None else PADY+IN(60)
    tf=textbox(slide,PADX,y,CW-PADX*2,IN(200))
    col=color or (WHITE if dark else FOREST)
    set_text(line0(tf),txt,PT(size),col,DISP,True)
    return y

def accent(slide,y,color=GREEN):
    rect(slide,PADX,y,IN(96),IN(8),fill=color)

def subhead(slide,txt,dark,y):
    tf=textbox(slide,PADX,y,IN(1300),IN(120))
    set_text(line0(tf),txt,PT(30),BODY_DARK if dark else META_LIGHT,BODY,False)

def footer(slide,dark):
    # rule
    rect(slide,PADX,CH-IN(190),CW-PADX*2,Emu(9525),fill=(RGBColor(0x2A,0x4A,0x4C) if dark else RGBColor(0xDD,0xE3,0xDD)))
    # logo plate
    plate=rect(slide,PADX,CH-IN(168),IN(300),IN(112),fill=WHITE)
    try: slide.shapes.add_picture(LOGO,PADX+IN(22),CH-IN(150),height=IN(76))
    except Exception: pass
    tf=textbox(slide,CW-PADX-IN(1200),CH-IN(160),IN(1200),IN(100),MSO_ANCHOR.MIDDLE)
    p=line0(tf); set_text(p,COMPANY_L1,PT(22),META_DARK if dark else META_LIGHT,BODY,False,PP_ALIGN.RIGHT)
    p2=para(tf,0,0,PP_ALIGN.RIGHT); set_text(p2,COMPANY_L2,PT(22),META_DARK if dark else META_LIGHT,BODY,False,PP_ALIGN.RIGHT)

def std_header(slide,d,dark):
    y=PADY
    if d.get('eyebrow'):
        eyebrow(slide,d['eyebrow'],dark); y=PADY+IN(58)
    if d.get('headline'):
        hy=headline(slide,d['headline'],dark,y=y,size=(60 if len(d['headline'])>34 else 72))
        y=hy+IN(150 if len(d['headline'])>34 else 110)
        accent(slide,y,d.get('accentcolor',GREEN)); y+=IN(34)
    if d.get('subhead'):
        subhead(slide,d['subhead'],dark,y); y+=IN(110)
    return y

# ---------------------------------------------------------------------------
# CONTENT (mirrors deck/content/modals.js + slides.js)
# ---------------------------------------------------------------------------
def S(head,items,tone=None,note=None): return dict(head=head,items=items,tone=tone,note=note)

MODALS={
 'myth-20-down':dict(eyebrow='The myth',title='You need 20% down',sections=[
   S('The reality',['Conventional goes as low as 3%, FHA 3.5%, VA and USDA 0%','The median first-time buyer puts down far less than 20%']),
   S('Pros of 20% down',['No monthly mortgage insurance','Lower payment and more instant equity'],'pros'),
   S('Cons',['Years of saving while rent builds no equity','Drains the cash reserve you want after closing'],'cons'),
   S(None,None,note='Mortgage insurance on a conventional loan is removable.')]),
 'myth-fha-first':dict(eyebrow='The myth',title='FHA is only for first-time buyers',sections=[
   S('The reality',['FHA has no first-time buyer requirement — anyone eligible can use it','It’s a primary-residence program, not a first-timer program']),
   S('Pros',['More forgiving on credit and past events','Higher DTI tolerance · assumable'],'pros'),
   S('Cons',['MI often lasts the life of the loan under 10% down','Upfront premium added to the balance'],'cons')]),
 'myth-perfect-credit':dict(eyebrow='The myth',title='You need perfect credit',sections=[
   S('The reality',['Pricing moves in tiers — near an edge matters, mid-tier often doesn’t','FHA and VA reach meaningfully lower than conventional']),
   S('Pros of a higher score',['Better pricing, lower MI, more program options'],'pros'),
   S('Cons of waiting for "perfect"',['Prices and rents keep moving while you chase a number'],'cons'),
   S(None,None,note='<strong>Utilization</strong> is ~30% of the score and can move in one cycle. Don’t close old cards.')]),
 'myth-renting':dict(eyebrow='The myth',title='Renting is cheaper',sections=[
   S('The reality',['Rent generally rises; a fixed principal & interest payment doesn’t','Ownership builds equity — rent doesn’t']),
   S('Pros of renting',['Flexibility, no maintenance, no transaction costs'],'pros'),
   S('Cons of renting',['No equity, no fixed housing cost, no control over increases'],'cons'),
   S(None,None,note='Short horizon, renting can win; long horizon, ownership usually does.')]),
 'myth-lowest-rate':dict(eyebrow='The myth',title='The lowest rate is the best deal',sections=[
   S('The reality',['A lower rate is often bought with points paid at closing','The lowest rate and the lowest cost are never the same loan']),
   S('Pros of a lower rate',['Lower payment and less interest over a long hold','The benefit of a lower Rate is lower payment.'],'pros'),
   S('Cons',['Points are cash today — gone if you sell or refinance early','Chasing rate can hide high lender fees'],'cons'),
   S(None,None,note='Ask for <strong>rate AND total cost</strong>, and compare Loan Estimates from the same day.')]),
 'prog-conventional':dict(eyebrow='Loan program',title='Conventional',compliance=True,sections=[
   S('Who it fits',['Solid credit, steady documented income · as low as 3% down']),
   S('Pros',['MI comes off','Best pricing at strong credit','Widest property flexibility'],'pros'),
   S('Cons',['Least forgiving on credit','Less tolerant of recent derogatory history'],'cons')]),
 'prog-fha':dict(eyebrow='Loan program',title='FHA',compliance=True,sections=[
   S('Who it fits',['Rebuilding credit or a higher debt load · 3.5% down']),
   S('Pros',['Lower credit thresholds','Higher DTI tolerance','Assumable'],'pros'),
   S('Cons',['MI usually lasts the life of the loan','Upfront premium added to the balance'],'cons')]),
 'prog-va':dict(eyebrow='Loan program',title='VA',compliance=True,sections=[
   S('Who it fits',['Eligible service members, veterans, and certain spouses · 0% down']),
   S('Pros',['No down payment','No monthly MI','Reusable — not once per lifetime'],'pros'),
   S('Cons',['A funding fee applies (financeable; waived for some)','Stricter property standards'],'cons'),
   S(None,None,note='If anyone in the household served — including a spouse — <strong>ask.</strong> It usually counts.')]),
 'prog-usda':dict(eyebrow='Loan program',title='USDA',compliance=True,sections=[
   S('Who it fits',['Eligible area, under the household income limit · 0% down']),
   S('Pros',['No down payment','Annual fee often lower than FHA'],'pros'),
   S('Cons',['Geography decides — check the map','Income limits count everyone in the home'],'cons')]),
 'src-savings':dict(eyebrow='Where it comes from',title='Checking & Savings',sections=[
   S('The rules',['Must be seasoned — sitting in your account, documented','No undocumented cash. Cash deposits without a trail can’t be used','Large non-payroll deposits must be sourced and explained'],note='Move money <strong>before</strong> you apply, not during.')]),
 'src-investments':dict(eyebrow='Where it comes from',title='Investment Accounts',sections=[
   S('The rules',['Brokerage, stocks, and mutual funds count — a recent statement documents them','To use the money you sell the assets, then show the sale and the deposit','For reserves, underwriting may count only part of the balance, not the full amount'],note='You don’t have to sell to <strong>prove</strong> reserves — only to <strong>spend</strong> them.')]),
 'src-retirement':dict(eyebrow='Where it comes from',title='Retirement Accounts',sections=[
   S('The rules',['Vested 401(k) and IRA balances count — a statement documents them','To use the funds, document the withdrawal or loan and the money landing in your account','A 401(k) loan payment can count in your DTI; a withdrawal may bring taxes or penalties'],note='Weigh the tax hit with your advisor <strong>before</strong> you pull from retirement.')]),
 'src-gift':dict(eyebrow='Where it comes from',title='Gift Funds',sections=[
   S('The rules',['Gift letter required — stating it is not a loan','Generally must come from family','The donor’s funds must be sourced, and the transfer documented'],note='Tell your lender <strong>before</strong> the money moves.')]),
 'src-seller':dict(eyebrow='Where it comes from',title='Seller Credits',sections=[
   S('The rules',['The seller contributes toward your closing costs','Concession caps apply by program and down payment','Cannot pay your down payment or exceed your actual costs'],note='You have to <strong>negotiate</strong> for it — ask your agent to ask.')]),
 'src-buydown':dict(eyebrow='Where it comes from',title='2-1 Buydowns',compliance=True,sections=[
   S('How it works',['A temporary payment reduction for the first two years','Usually paid by the seller or builder, not you']),
   S('The catch',['You qualify at the real rate — full payment in year three','A plan to refinance first depends on rates nobody controls'],'cons')]),
 'src-assistance':dict(eyebrow='Where it comes from',title='Down Payment Assistance',sections=[
   S('What it is',['Grants or second mortgages that reduce cash to close']),
   S('The trade-offs',['Complicates the process — extra underwriting and timelines','Often makes your offer less attractive to sellers','Usually a higher rate, caps, and fine print'],'cons'),
   S(None,None,note='Sometimes it clearly wins. Run it against the higher rate before deciding.')]),
 'role-lender':dict(eyebrow='Your team',title='Lender / Loan Officer',sections=[
   S('What they do',['Structures your financing and guides you from pre-approval to closing','Packages the file and advocates for it'],note='We don’t approve the loan — <strong>the underwriter does.</strong>')]),
 'role-realtor':dict(eyebrow='Your team',title='Realtor',sections=[
   S('What they do',['Finds homes, writes and negotiates your offer, manages deadlines'],note='The listing agent works for the <strong>seller</strong>. Yours works for you — ask how they’re paid.')]),
 'role-seller':dict(eyebrow='Your team',title='Seller',sections=[
   S('What they do',['Owns the home; their motivation shapes what’s negotiable','You almost never deal with them directly — it runs through agents'],note='Sometimes <strong>timing or certainty</strong> matters to them more than top price.')]),
 'role-underwriter':dict(eyebrow='Your team',title='Underwriter',sections=[
   S('What they do',['Makes the actual decision against program guidelines','You’ll never deal with them directly'],note='Not looking for a reason to say no — documenting a reason to say yes.')]),
 'role-title':dict(eyebrow='Your team',title='Title Company',sections=[
   S('What they do',['Researches ownership, clears liens, issues title insurance, holds funds','Neutral to both sides'],note='You can often <strong>shop</strong> for them — most people never do.')]),
 'role-closer':dict(eyebrow='Your team',title='Closer',sections=[
   S('What they do',['Runs the signing, walks you through documents, moves the funds'],note='Budget an hour. Read what you sign and ask about anything unclear.')]),
 'role-insurance':dict(eyebrow='Your team',title='Insurance Agent',sections=[
   S('What they do',['Places your homeowners policy, required before closing'],note='<strong>Start early.</strong> Premiums vary widely and a high one can affect qualifying.')]),
 'role-processor':dict(eyebrow='Your team',title='Processor',sections=[
   S('What they do',['Collects your documents, orders appraisal and title, preps the file'],note='Being asked twice isn’t a lost document — <strong>send it again same-day.</strong>')]),
 'assume-preapproval':dict(eyebrow='Small assumption',title='"I’ll get pre-approved once I find a house"',sections=[
   S('What actually happens',['The house goes under contract while you’re getting a letter']),
   S('Do instead',['Pre-approval first — it’s free and nothing is committed'],'pros')]),
 'assume-my-bank':dict(eyebrow='Small assumption',title='"I’ll just use my bank"',sections=[
   S('What actually happens',['A bank offers its own products; fees and pricing can run thousands higher at the same rate']),
   S('Do instead',['Get more than one Loan Estimate the same day and compare fees','A broker shops many lenders for you — and it counts as one credit pull'],'pros')]),
 'assume-closing-costs':dict(eyebrow='Small assumption',title='"Closing costs are what I need at closing"',sections=[
   S('What actually happens',['You arrive short — cash to close also has down payment, prepaids, escrows, reserves']),
   S('Do instead',['Plan around Cash to Close on page one, not closing costs'],'pros')]),
 'assume-approved':dict(eyebrow='Small assumption',title='"I’m approved, so I’m done"',sections=[
   S('What actually happens',['Lenders re-pull credit before closing — new debt can break the file']),
   S('Do instead',['Change nothing until you have keys. If something must change, call first'],'pros')]),
 'assume-wait-20':dict(eyebrow='Small assumption',title='"I’ll wait until I have twenty percent"',sections=[
   S('What actually happens',['The target rises while you save — and you pay rent the whole time']),
   S('Do instead',['Run the comparison. Waiting should be a decision, not a default'],'pros')]),
 'assume-builder':dict(eyebrow='Small assumption',title='"The builder’s incentive is free money"',sections=[
   S('What actually happens',['It may be funded by a higher rate, higher fees, or a price set to pay for it']),
   S('Do instead',['Take the builder’s Loan Estimate, get one from outside, and compare fees'],'pros')]),
}
print("MODALS:",len(MODALS))

# ---- main slides -----------------------------------------------------------
def C(modal,title,meta=None,stat=None): return dict(modal=modal,title=title,meta=meta,stat=stat)

SLIDES=[
 dict(id='opening',layout='opening',bg='dark',footer=True,eyebrow="The Homebuyer's Playbook",headline='Buying your first home, explained.'),
 dict(id='myths',layout='grid',bg='white',footer=True,cols=3,headline='Myths',subhead='The stories that keep good buyers renting. Tap any one.',cards=[
   C('myth-20-down','You need 20% down'),C('myth-fha-first','FHA is only for first-time buyers'),C('myth-perfect-credit','You need perfect credit'),
   C('myth-renting','Renting is cheaper'),C('myth-lowest-rate','The lowest rate is the best deal')]),
 dict(id='budget-rent-buy',layout='compare',bg='dark',eyebrow='Budget',
   left=dict(label='Renting',items=['Rent in Denver has trended up, year after year','Flexible — but every payment builds someone else’s equity','The longer you wait, the more expensive buying becomes.']),
   right=dict(label='Buying',items=['A fixed payment while home values generally appreciate','Equity builds two ways — paying down principal and appreciation','Most people step up — the first home isn’t the last','You can personalize the home to fit your style and needs.']),
   callout='Waiting isn’t neutral. It usually costs something.'),
 dict(id='budget-payment',layout='payment',bg='white',footer=True,eyebrow='Budget',headline="What's actually in the payment",
   subhead='On a fixed-rate loan — what’s locked, and what can still move.',
   fixed=['Your interest rate','Your principal & interest','Your loan term'],
   moves=['Property taxes','Insurance premiums','Mortgage insurance (it can come off)','HOA dues','A special assessment'],
   points=['Escrow collects taxes and insurance monthly and pays them for you — not a fee, your money held early',
     'Amortization: early payments are mostly interest; it shifts to principal over years',
     'Refinancing resets the clock — a new 30-year loan starts over at the interest-heavy beginning']),
 dict(id='budget-comfort',layout='points',bg='dark',eyebrow='Budget',headline='Keep the payment in your comfort zone',
   subhead='The number you can qualify for and the number you can live with are rarely the same.',points=[
   "Don’t become house poor — the guidelines ask if you’ll repay, not if you’ll be okay",
   'Protect quality of life — leave room for disposable income and emergencies',
   'The payment can change after you close']),
 dict(id='programs',layout='grid',bg='white',footer=True,cols=2,headline='Most Common Loan Programs',subhead='Four doors, different people. Tap any program for pros and cons.',cards=[
   C('prog-conventional','Conventional','The default for most buyers','3% down'),C('prog-usda','USDA','Eligible areas, income limits','0% down'),
   C('prog-fha','FHA','Built for credit flexibility','3.5% down'),C('prog-va','VA','If you served, check this first','0% down')]),
 dict(id='cash-sources',layout='grid',bg='dark',footer=True,cols=3,headline='Cash to close: where it comes from',subhead='It doesn’t all have to sit in one account. Tap any source for the rules.',cards=[
   C('src-savings','Checking & Savings'),C('src-investments','Investment Accounts'),C('src-retirement','Retirement Accounts'),C('src-gift','Gift Funds'),C('src-seller','Seller Credits'),C('src-buydown','2-1 Buydowns'),C('src-assistance','Down Payment Assistance')]),
 dict(id='cash-makeup',layout='cashmakeup',bg='white',footer=True,numbers=True,eyebrow='Cash to close',headline='Closing costs are not cash to close',
   subhead='Closing costs are the fees to set up the loan. Cash to close is the total amount you bring to closing — and it’s the bigger number.',
   canpay=['Allowable closing costs','Prepaid items','Escrow deposits','Points to lower the rate'],
   cannotpay=['Your down payment','More than your actual costs','Cash back to you'],
   points=['You can use points to buy down the rate — or credits to cover closing costs',
     'Prepaids (first insurance, some taxes) are your own money paid early, not a fee',
     'Reserves are funds you keep after closing — some programs require them']),
 dict(id='process-players',layout='grid',bg='white',footer=True,cols=4,dense=True,headline='Meet the players',subhead='Eight people — and they don’t all work for you. Tap anyone.',cards=[
   C('role-lender','Lender'),C('role-realtor','Realtor'),C('role-seller','Seller'),C('role-processor','Processor'),
   C('role-underwriter','Underwriter'),C('role-title','Title'),C('role-closer','Closer'),C('role-insurance','Insurance')]),
 dict(id='process-steps',layout='stepper',bg='dark',footer=True,eyebrow='The process',headline='From first call to keys',steps=[
   ('Pre-Approval','1–2 days'),('Application','Same day'),('Processing','1–2 weeks'),('Underwriting','Days–1 week'),
   ('Conditional Approval','A few days'),('Clear to Close','Days before'),('Closed','About an hour'),('Funded','Same / next day')]),
 dict(id='mistakes-donts',layout='markers',bg='white',footer=True,tone='dont',cols=2,eyebrow='Mistakes to avoid',headline="Don't",
   subhead='From application to closing, your financial life is a museum exhibit — look, don’t touch.',items=[
   'Apply for or take out any new loans','Pay off collections or debt unless your lender says to','Consolidate debt or close credit cards',
   'Dispute items on your credit report','Make large deposits other than payroll — no cash','Make large purchases with credit or debit',
   'Move money between accounts','Change jobs without talking to your loan officer','Quit or give notice before your loan closes',
   'Use cash for earnest money — check or wire only']),
 dict(id='mistakes-dos',layout='markers',bg='mist',footer=True,tone='do',cols=1,eyebrow='Mistakes to avoid',headline='Do',subhead='The short list that keeps a file moving.',items=[
   'Make all your payments on time','Tell your lender if you haven’t filed last year’s taxes, or owe on them',
   'Tell your lender if anything about your job is changing','Keep homeowners insurance at or under the Loan Estimate amount',
   'Return everything your lender asks for quickly']),
 dict(id='mistakes-rate',layout='points',bg='dark',eyebrow='Mistakes to avoid',headline='Don’t assume the lowest rate wins',points=[
   'Breakeven: a rate bought with points only pays off if you keep the loan long enough',
   'Points are a bet on how long you keep this exact loan — not this house',
   'Refinancing ends the loan early and could restart amortization',
   'Real wealth comes from building equity — payments and time'],
   callout='Equity is built by payments and time — not the third decimal of your rate.'),
 dict(id='mistakes-assumptions',layout='grid',bg='white',footer=True,cols=3,quote=True,eyebrow='Mistakes to avoid',headline='The most expensive mistakes start as small assumptions',subhead='Tap any one for what actually happens.',cards=[
   C('assume-preapproval','“I’ll get pre-approved once I find a house.”'),C('assume-my-bank','“I’ll just use my bank.”'),
   C('assume-closing-costs','“Closing costs are what I need at closing.”'),C('assume-approved','“I’m approved, so I’m done.”'),
   C('assume-wait-20','“I’ll wait until I have twenty percent.”'),C('assume-builder','“The builder’s incentive is free money.”')]),
 dict(id='questions',layout='questions',bg='dark',footer=True,eyebrow='Questions',headline='The questions everyone asks',items=[
   ('How much money do I actually need?','Less than most think — 3% conventional, 3.5% FHA, 0% VA/USDA — plus closing costs and an emergency fund you don’t touch.'),
   ('What credit score do I need?','No single number. It depends on the program, and you’re probably closer than you think.'),
   ('How long does it take?','Contract to keys is typically 30–45 days. The biggest variable is how fast you return documents.'),
   ('Should I lock my rate?','Generally once you’re under contract. A lock holds your rate while the loan is in process — ask how long it lasts and what happens if closing runs past it.'),
   ('What happens after pre-approval?','You shop with a real budget, make an offer, and it becomes a formal application — Loan Estimate within three business days.')]),
 dict(id='wrap',layout='wrap',bg='dark',footer=True,headline='Let’s build your plan.'),
]
print("SLIDES:",len(SLIDES))

# ---------------------------------------------------------------------------
# RENDERERS
# ---------------------------------------------------------------------------
main_slides=[]; card_links=[]; modal_parent={}

def new_slide(bgcolor):
    s=prs.slides.add_slide(BLANK); bg(s,bgcolor); return s

def bullets_box(slide,l,t,w,h,items,size,color,gap=8,bcolor=GREEN,marker=True):
    tf=textbox(slide,l,t,w,h)
    for i,it in enumerate(items):
        p=line0(tf) if i==0 else para(tf,space_after=gap)
        if marker:
            rr=p.add_run(); rr.text='■  '; rr.font.name=BODY; rr.font.size=size; rr.font.color.rgb=bcolor
        add_rich(p,it,size,color)
    return tf

def add_card(slide,x,y,w,h,card,dark,quote=False,border=GREEN):
    rect(slide,x,y,IN(8),h,fill=border)                       # left accent
    body=rect(slide,x+IN(8),y,w-IN(8),h,fill=(RGBColor(0x1B,0x40,0x42) if dark else WHITE),
              line=(None if dark else RGBColor(0xE2,0xE7,0xE2)))
    tf=body.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=IN(30); tf.margin_right=IN(30); tf.margin_top=IN(28); tf.margin_bottom=IN(20)
    p=line0(tf)
    if quote: set_text(p,card['title'],PT(30),(BODY_DARK if dark else CHARCOAL),BODY,False)
    else:     set_text(p,card['title'],PT(34 if not card.get('meta') else 34),(WHITE if dark else FOREST),DISP,True)
    if card.get('meta'):
        pm=para(tf,4,4); set_text(pm,card['meta'],PT(24),(BODY_DARK if dark else CHARCOAL),BODY,False)
    if card.get('stat'):
        ps=para(tf,10,0); set_text(ps,card['stat'],PT(40),(GREEN if dark else SAGE),DISP,True)
    # + cue
    plus=textbox(slide,x+w-IN(56),y+IN(18),IN(40),IN(40),MSO_ANCHOR.MIDDLE)
    set_text(line0(plus),'+',PT(40),(WHITE if dark else FOREST),DISP,True,PP_ALIGN.CENTER)
    return body

def L_opening(s,d,dark):
    try: s.shapes.add_picture(PORTRAIT,IN(1180),0,height=CH)
    except Exception: rect(s,IN(1180),0,CW-IN(1180),CH,fill=FOREST_MID)
    rect(s,IN(1180),0,CW-IN(1180),CH,fill=None)
    eyebrow(s,d['eyebrow'],dark)
    headline(s,d['headline'],dark,y=PADY+IN(70),size=76)
    accent(s,PADY+IN(430))
    tf=textbox(s,PADX,PADY+IN(500),IN(1000),IN(300))
    set_text(line0(tf),SETH['name'],PT(46),WHITE,DISP,True)
    p=para(tf,2,18); set_text(p,SETH['title'],PT(28),BODY_DARK,BODY,False)
    for lbl,val in [('CALL',SETH['phone']),('EMAIL',SETH['email']),('WEB','msfg.us')]:
        pc=para(tf,4,4); r=pc.add_run(); r.text=lbl+'   '; r.font.name=DISP; r.font.bold=True; r.font.size=PT(22); r.font.color.rgb=GREEN
        r2=pc.add_run(); r2.text=val; r2.font.name=BODY; r2.font.size=PT(30); r2.font.color.rgb=WHITE

def L_grid(s,d,dark):
    y=std_header(s,d,dark)
    cols=d.get('cols',3); dense=d.get('dense'); quote=d.get('quote')
    n=len(d['cards']); gap=IN(28)
    gridw=CW-PADX*2
    cw=(gridw-gap*(cols-1))/cols
    ch=IN(150 if dense else (150 if quote else 175))
    rows=(n+cols-1)//cols
    top=y+IN(10)
    border = RED if d['id']=='mistakes-donts' else GREEN
    for i,card in enumerate(d['cards']):
        r=i//cols; c=i%cols
        x=PADX+c*(cw+gap); yy=top+r*(ch+gap)
        shape=add_card(s,x,yy,cw,ch,card,dark,quote=quote,border=border)
        modal_parent[card['modal']]=s; card_links.append((shape,card['modal']))

def L_points(s,d,dark):
    y=std_header(s,d,dark)
    bullets_box(s,PADX,y+IN(20),IN(1500),IN(500),d['points'],PT(30),(BODY_DARK if dark else CHARCOAL),gap=12)
    if d.get('callout'):
        coy = CH-IN(420) if d.get('compare') else CH-IN(330)
        co=rect(s,PADX,coy,IN(1500),IN(84),fill=GREEN)
        tf=co.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.margin_left=IN(40)
        set_text(line0(tf),d['callout'],PT(32),FOREST,DISP,True)
    if d.get('compare'):
        cb=rect(s,PADX,CH-IN(300),IN(820),IN(84),fill=GREEN); tfc=cb.text_frame; tfc.vertical_anchor=MSO_ANCHOR.MIDDLE
        set_text(line0(tfc),'Compare loans:  Conventional · FHA · VA  ›',PT(24),FOREST,DISP,True,PP_ALIGN.CENTER)
        modal_parent[d['compare']]=s; card_links.append((cb,d['compare']))

def L_compare(s,d,dark):
    half=CW/2
    rect(s,0,0,half,CH,fill=WHITE); rect(s,half,0,half,CH,fill=FOREST)
    # left
    tf=textbox(s,IN(96),IN(140),half-IN(180),IN(500))
    set_text(line0(tf),d['eyebrow'].upper(),PT(23),SAGE,DISP,True)
    pl=para(tf,20,24); set_text(pl,d['left']['label'],PT(54),FOREST,DISP,True)
    for it in d['left']['items']:
        p=para(tf,0,18); rr=p.add_run(); rr.text='■  '; rr.font.color.rgb=GREEN; rr.font.size=PT(30); rr.font.name=BODY
        add_rich(p,it,PT(30),CHARCOAL)
    # right
    tf2=textbox(s,half+IN(84),IN(220),half-IN(180),IN(500))
    set_text(line0(tf2),d['right']['label'],PT(54),WHITE,DISP,True)
    for it in d['right']['items']:
        p=para(tf2,0,18); rr=p.add_run(); rr.text='■  '; rr.font.color.rgb=GREEN; rr.font.size=PT(30); rr.font.name=BODY
        add_rich(p,it,PT(30),BODY_DARK)
    if d.get('callout'):
        cw=IN(1000); co=rect(s,(CW-cw)/2,CH-IN(210),cw,IN(90),fill=GREEN)
        tf3=co.text_frame; tf3.vertical_anchor=MSO_ANCHOR.MIDDLE
        set_text(line0(tf3),d['callout'],PT(32),FOREST,DISP,True,PP_ALIGN.CENTER)

def L_payment(s,d,dark):
    # compact header so the bar + columns + explainer row all fit above the footer
    eyebrow(s,d['eyebrow'],dark)
    headline(s,d['headline'],dark,y=PADY+IN(46),size=56)
    accent(s,PADY+IN(150))
    subhead(s,d['subhead'],dark,PADY+IN(180))
    top=PADY+IN(250)
    # banded bar
    segs=[('Principal',210),('Interest',300),('Taxes',170),('Insurance',120),('Mortgage ins.',110),('HOA',90)]
    total=sum(w for _,w in segs); barw=IN(1580); x=PADX; barh=IN(72)
    for i,(lbl,w) in enumerate(segs):
        sw=barw*w/total; rect(s,x,top,sw,barh,fill=PAY[i])
        lt=textbox(s,x,top+barh+IN(6),sw,IN(40)); set_text(line0(lt),lbl,PT(19),CHARCOAL,DISP,True,PP_ALIGN.CENTER)
        x+=sw
    cy=top+barh+IN(66)
    # two columns
    colw=IN(740)
    for ci,(head,items,bc) in enumerate([("What's locked",d['fixed'],SAGE),('What can still move',d['moves'],GREEN)]):
        cx=PADX+ci*(colw+IN(60))
        ht=textbox(s,cx,cy,colw,IN(46)); set_text(line0(ht),head,PT(24),FOREST,DISP,True)
        rect(s,cx,cy+IN(44),colw,Emu(28575),fill=bc)
        bullets_box(s,cx,cy+IN(60),colw,IN(200),items,PT(22),CHARCOAL,gap=5)
    # explainer row, positioned below the taller (5-item) column
    py=cy+IN(240)
    pw=(CW-PADX*2-IN(68))/3
    for i,pt in enumerate(d['points']):
        px=PADX+i*(pw+IN(34))
        bullets_box(s,px,py,pw,IN(130),[pt],PT(21),CHARCOAL,gap=3)

def L_cashmakeup(s,d,dark):
    eyebrow(s,d['eyebrow'],dark)
    headline(s,d['headline'],dark,y=PADY+IN(46),size=56)
    accent(s,PADY+IN(150))
    subhead(s,d['subhead'],dark,PADY+IN(180))
    top=PADY+IN(250)
    colw=IN(760)
    for ci,(head,items,bc) in enumerate([('Credits CAN pay',d['canpay'],SAGE),('Credits CANNOT pay',d['cannotpay'],GREEN)]):
        cx=PADX+ci*(colw+IN(60))
        ht=textbox(s,cx,top,colw,IN(46)); set_text(line0(ht),head,PT(26),FOREST,DISP,True)
        rect(s,cx,top+IN(46),colw,Emu(28575),fill=bc)
        bullets_box(s,cx,top+IN(64),colw,IN(220),items,PT(24),CHARCOAL,gap=6)
    py=top+IN(250)
    bullets_box(s,PADX,py,IN(1560),IN(240),d['points'],PT(24),CHARCOAL,gap=10)

def L_stepper(s,d,dark):
    y=std_header(s,d,dark)
    steps=d['steps']; n=len(steps)
    x0=PADX; xw=CW-PADX*2; midy=CH/2-IN(20)
    seg=xw/(n-1) if n>1 else 0
    rect(s,x0,midy-Emu(19050),xw,Emu(38100),fill=RGBColor(0x2A,0x4A,0x4C))
    for i,(lbl,note) in enumerate(steps):
        cx=x0+i*seg; col=BANDS[min(4,round(i/(n-1)*4))]
        dia=IN(52); rect(s,cx-dia/2,midy-dia/2,dia,dia,fill=col,shape=MSO_SHAPE.OVAL)
        dot=IN(18); rect(s,cx-dot/2,midy-dot/2,dot,dot,fill=FOREST,shape=MSO_SHAPE.OVAL)
        lt=textbox(s,cx-seg/2,midy+IN(60),seg,IN(70)); set_text(line0(lt),lbl,PT(24 if len(lbl)<15 else 20),WHITE,DISP,True,PP_ALIGN.CENTER)
        nt=textbox(s,cx-seg/2,midy+IN(130),seg,IN(50)); set_text(line0(nt),note,PT(21),META_DARK,BODY,False,PP_ALIGN.CENTER)

def L_markers(s,d,dark):
    y=std_header(s,d,dark)
    tone=d['tone']; items=d['items']
    color = RED if tone=='dont' else DO_GREEN
    glyph = '✕' if tone=='dont' else '+'
    cols=d.get('cols',1); per=(len(items)+cols-1)//cols
    colw=(CW-PADX*2-IN(64))/cols
    top=y+IN(20)
    for i,it in enumerate(items):
        c=i//per; r=i%per; x=PADX+c*(colw+IN(64)); yy=top+r*IN(78)
        sq=rect(s,x,yy,IN(46),IN(46),fill=color); tfq=sq.text_frame; tfq.vertical_anchor=MSO_ANCHOR.MIDDLE
        set_text(line0(tfq),glyph,PT(26),WHITE,DISP,True,PP_ALIGN.CENTER)
        tf=textbox(s,x+IN(66),yy,colw-IN(80),IN(70),MSO_ANCHOR.MIDDLE); add_rich(line0(tf),it,PT(26 if cols==2 else 30),(FOREST if not dark else WHITE) if False else CHARCOAL)

def L_questions(s,d,dark):
    # compact header (no subhead) + tighter rows so all five clear the footer
    eyebrow(s,d['eyebrow'],dark)
    headline(s,d['headline'],dark,y=PADY+IN(46),size=60)
    accent(s,PADY+IN(158))
    top=PADY+IN(208)
    for i,(q,a) in enumerate(d['items']):
        yy=top+i*IN(116)
        sq=rect(s,PADX,yy,IN(50),IN(50),fill=GREEN); tq=sq.text_frame; tq.vertical_anchor=MSO_ANCHOR.MIDDLE
        set_text(line0(tq),str(i+1),PT(26),FOREST,DISP,True,PP_ALIGN.CENTER)
        tf=textbox(s,PADX+IN(78),yy,IN(1520),IN(120))
        set_text(line0(tf),q,PT(30),WHITE,DISP,True)
        pa=para(tf,3,0); add_rich(pa,a,PT(22),BODY_DARK)

def L_wrap(s,d,dark):
    headline(s,d['headline'],dark,y=PADY+IN(30),size=76)
    accent(s,PADY+IN(210))
    tf=textbox(s,PADX,PADY+IN(280),IN(1000),IN(400))
    set_text(line0(tf),SETH['name'],PT(42),WHITE,DISP,True)
    p=para(tf,2,18); set_text(p,SETH['title'],PT(26),BODY_DARK,BODY,False)
    for lbl,val in [('CALL',SETH['phone']),('EMAIL',SETH['email']),('WEB','msfg.us')]:
        pc=para(tf,4,4); r=pc.add_run(); r.text=lbl+'   '; r.font.name=DISP; r.font.bold=True; r.font.size=PT(22); r.font.color.rgb=GREEN
        r2=pc.add_run(); r2.text=val; r2.font.name=BODY; r2.font.size=PT(30); r2.font.color.rgb=WHITE
    # Apply Now — real hyperlink. Schedule button omitted until the booking link exists.
    b1=rect(s,PADX,CH-IN(400),IN(430),IN(90),fill=GREEN); tb=b1.text_frame; tb.vertical_anchor=MSO_ANCHOR.MIDDLE
    set_text(line0(tb),'APPLY NOW',PT(24),FOREST,DISP,True,PP_ALIGN.CENTER)
    try: b1.click_action.hyperlink.address=APPLY_URL
    except Exception: pass
    try: s.shapes.add_picture(QR,CW-IN(560),PADY+IN(120),height=IN(420))
    except Exception: pass

LAYOUTS=dict(opening=L_opening,grid=L_grid,points=L_points,compare=L_compare,payment=L_payment,
             cashmakeup=L_cashmakeup,stepper=L_stepper,markers=L_markers,questions=L_questions,wrap=L_wrap)

# ---- build main slides ----
for d in SLIDES:
    dark = d['bg']=='dark'
    s=new_slide(FOREST if dark else (MIST if d['bg']=='mist' else WHITE))
    if d['id']=='mistakes-donts': d['accentcolor']=RED
    if d['id']=='mistakes-dos': d['accentcolor']=DO_GREEN
    if d['id']=='mistakes-donts': d['headline_color']=RED
    LAYOUTS[d['layout']](s,d,dark)
    if d.get('footer'): footer(s,dark)
    main_slides.append(s)

# ---- popout slides ----
def render_popout(mid,m):
    s=new_slide(WHITE)
    # header band
    dark=True
    band=rect(s,0,0,CW,IN(300),fill=FOREST)
    tf=textbox(s,PADX,IN(56),CW-PADX*2-IN(200),IN(220))
    if m.get('eyebrow'): set_text(line0(tf),m['eyebrow'].upper(),PT(23),GREEN,DISP,True); ph=para(tf,10,0)
    else: ph=line0(tf)
    set_text(ph,m['title'],PT(50 if len(m['title'])<40 else 40),WHITE,DISP,True)
    rect(s,PADX,IN(250),IN(96),IN(8),fill=GREEN)
    # body
    by=IN(360)
    if m.get('intro'):
        it=textbox(s,PADX,by,CW-PADX*2,IN(140)); set_text(line0(it),m['intro'],PT(25),CHARCOAL,BODY,False); by=IN(480)
    if m.get('table'):
        t=m['table']; cols=len(t['columns'])+1; rows=len(t['rows'])+1
        gt=s.shapes.add_table(rows,cols,PADX,by,CW-PADX*2,IN(52*rows)).table
        gt.first_row=True; gt.horz_banding=True
        hdr=['']+t['columns']
        for c,txt in enumerate(hdr):
            cell=gt.cell(0,c); cell.fill.solid(); cell.fill.fore_color.rgb=FOREST
            pp=cell.text_frame.paragraphs[0]; set_text(pp,txt,PT(24),WHITE,DISP,True)
        for r,(label,cells) in enumerate(t['rows'],start=1):
            lc=gt.cell(r,0); lc.fill.solid(); lc.fill.fore_color.rgb=MIST
            set_text(lc.text_frame.paragraphs[0],label,PT(20),FOREST,DISP,True)
            for c,val in enumerate(cells,start=1):
                cc=gt.cell(r,c); cc.fill.solid(); cc.fill.fore_color.rgb=WHITE
                set_text(cc.text_frame.paragraphs[0],re.sub('<[^>]+>','',val),PT(20),CHARCOAL,BODY,False)
        by = by + IN(54*rows) + IN(20)
    tf2=textbox(s,PADX,by,CW-PADX*2,IN(700)-(by-IN(360)))
    first=True
    for sec in m['sections']:
        if sec['head']:
            ph=line0(tf2) if first else para(tf2,18,10); first=False
            hc = SAGE if sec['tone']=='pros' else (CHARCOAL if sec['tone']=='cons' else FOREST)
            pre = '+  ' if sec['tone']=='pros' else ('–  ' if sec['tone']=='cons' else '')
            set_text(ph,pre+sec['head'],PT(22),hc,DISP,True)
        for it in (sec['items'] or []):
            p=line0(tf2) if first else para(tf2,0,8); first=False
            bc = SAGE if sec['tone']=='pros' else (CHARCOAL if sec['tone']=='cons' else GREEN)
            rr=p.add_run(); rr.text='■  '; rr.font.color.rgb=bc; rr.font.size=PT(26); rr.font.name=BODY
            add_rich(p,it,PT(26),CHARCOAL)
        if sec['note']:
            p=para(tf2,14,8); first=False
            rr=p.add_run(); rr.text='  '; add_rich(p,sec['note'],PT(26),FOREST)
    if m.get('compliance'):
        txt=('Hypothetical illustration for education only. Not a quote, offer, or commitment to lend.'
             if m['compliance']=='hypothetical'
             else 'General guidelines only. Actual eligibility depends on full underwriting, credit, property, and lender overlays.')
        ct=textbox(s,PADX,CH-IN(210),CW-PADX*2,IN(80))
        set_text(line0(ct),txt,PT(21),META_LIGHT,BODY,False)
        line0(ct).runs[0].font.italic=True
    # back button
    back=rect(s,PADX,CH-IN(120),IN(240),IN(72),fill=FOREST); tb=back.text_frame; tb.vertical_anchor=MSO_ANCHOR.MIDDLE
    set_text(line0(tb),'‹  Back',PT(24),WHITE,DISP,True,PP_ALIGN.CENTER)
    return s,back

popout_slide={}; back_links=[]
for mid,m in MODALS.items():
    s,back=render_popout(mid,m); popout_slide[mid]=s; back_links.append((back,modal_parent.get(mid)))

# ---- wire hyperlinks ----
for shape,mid in card_links:
    if mid in popout_slide: shape.click_action.target_slide=popout_slide[mid]
for back,parent in back_links:
    if parent is not None: back.click_action.target_slide=parent

out=os.path.abspath(os.path.join(os.path.dirname(__file__),'..','Homebuyers-Playbook.pptx'))
prs.save(out)
print("saved:",out,"| slides:",len(prs.slides.__iter__.__self__._sldIdLst))
